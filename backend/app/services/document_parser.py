"""
Document parsing service.

Supports: txt, md, pdf, docx, pptx, xlsx/xls, and images (via ocr_service).

parse_document()             → List[str]  (backward compat, chunks only)
parse_document_with_media()  → (List[str], List[dict])  (chunks + image metadata)
"""

import asyncio
import base64
import json
import logging
import os
import re
import zipfile
from typing import List, Tuple, Dict, Any, Optional

from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif", ".webp"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}


# ── Text splitter ─────────────────────────────────────────────────────────────

def _split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""],
    )
    return splitter.split_text(text)


# ── Image helpers ─────────────────────────────────────────────────────────────

def _make_image_id(source_doc: str, seq_num: int, description: str = "") -> str:
    """Build a stable semantic image ID like ``IMG_运维知识100问电力_018``.

    Intentionally does NOT include the description: LLMs tend to truncate any
    appended description when citing markers (producing IDs that no longer
    match the stored row). The description lives in DocumentImage.description
    separately — we only need a stable, non-truncatable marker here.
    """
    doc_short = re.sub(r"[^\w\u4e00-\u9fff]", "", os.path.splitext(source_doc)[0])[:10]
    return f"IMG_{doc_short}_{seq_num:03d}"


async def _describe_image_with_vision(
    image_bytes: bytes,
    image_ext: str,
    vision_provider,
) -> str:
    """Generate a Chinese description of an image using an OCR/vision provider.

    `vision_provider` should be an OCRProvider instance (or any provider with
    encrypted_api_key/model_name fields). Returns "图片" on failure so indexing
    can continue.
    """
    from app.services import ocr_service
    return await ocr_service.describe_image_bytes(image_bytes, image_ext, vision_provider)


def _save_image(image_bytes: bytes, kb_id: str, doc_id: str, seq_num: int, ext: str) -> str:
    """Save image bytes to the uploads directory and return the local path."""
    from app.config import settings
    img_dir = os.path.join(settings.upload_dir, kb_id, "images")
    os.makedirs(img_dir, exist_ok=True)
    filename = f"{doc_id}_{seq_num:03d}{ext}"
    path = os.path.join(img_dir, filename)
    with open(path, "wb") as f:
        f.write(image_bytes)
    return path


# ── Table helpers ─────────────────────────────────────────────────────────────

def _table_to_markdown(table) -> str:
    """Convert a python-docx Table object to Markdown."""
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


# ── Plain text parsers ────────────────────────────────────────────────────────

def parse_txt(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return _split_text(text, chunk_size, chunk_overlap)


def parse_markdown(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return _split_text(text, chunk_size, chunk_overlap)


def parse_pdf(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    try:
        import PyPDF2
        text_parts = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text_parts.append(page.extract_text() or "")
        full_text = "\n".join(text_parts)
        return _split_text(full_text, chunk_size, chunk_overlap)
    except Exception as e:
        return [f"[PDF解析错误] {str(e)}"]


def parse_docx(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            parts.append(_table_to_markdown(table))
        full_text = "\n".join(parts)
        return _split_text(full_text, chunk_size, chunk_overlap)
    except Exception as e:
        return [f"[DOCX解析错误] {str(e)}"]


def parse_pptx(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        full_text = "\n".join(texts)
        return _split_text(full_text, chunk_size, chunk_overlap)
    except Exception as e:
        return [f"[PPTX解析错误] {str(e)}"]


def parse_excel(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    """Parse Excel (.xlsx/.xls) files. Each Sheet becomes an independent unit.
    Rows are grouped into chunks as Markdown tables preserving headers."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_chunks: List[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            # First row as header
            headers = [str(c) if c is not None else "" for c in rows[0]]
            data_rows = rows[1:]

            # Group rows into blocks of ~30 rows per chunk
            ROWS_PER_CHUNK = 30
            for start in range(0, max(1, len(data_rows)), ROWS_PER_CHUNK):
                block = data_rows[start: start + ROWS_PER_CHUNK]
                md_lines = [
                    f"## Sheet: {sheet_name}",
                    "",
                    "| " + " | ".join(headers) + " |",
                    "| " + " | ".join(["---"] * len(headers)) + " |",
                ]
                for row in block:
                    cells = [str(c) if c is not None else "" for c in row]
                    md_lines.append("| " + " | ".join(cells) + " |")
                all_chunks.extend(_split_text("\n".join(md_lines), chunk_size, chunk_overlap))

        wb.close()
        return all_chunks if all_chunks else ["[Excel文件为空]"]
    except ImportError:
        return ["[缺少 openpyxl 依赖，请执行: pip install openpyxl]"]
    except Exception as e:
        logger.exception("Excel解析失败: %s", file_path)
        return [f"[Excel解析错误] {str(e)}"]


def parse_document(file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    """Parse text-based documents → List[str] chunks. Backward-compatible."""
    ext = os.path.splitext(file_path)[1].lower()
    parsers = {
        ".txt": parse_txt,
        ".md": parse_markdown,
        ".pdf": parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".xlsx": parse_excel,
        ".xls": parse_excel,
    }
    parser = parsers.get(ext, parse_txt)
    return parser(file_path, chunk_size, chunk_overlap)


# ── Media-aware parsers ───────────────────────────────────────────────────────

async def _parse_docx_with_media(
    file_path: str,
    doc_id: str,
    kb_id: str,
    chunk_size: int,
    chunk_overlap: int,
    vision_provider,
    source_doc: str,
) -> Tuple[List[str], List[Dict[str, Any]], List[List[str]]]:
    """Parse DOCX with positional image extraction and standalone table chunks.

    Walks the document body in order. For each paragraph, extracts inline drawings
    and inserts ``[IMG_xxx]`` markers at their positions. Tables become standalone
    Markdown chunks.

    Returns:
        text_chunks:      list of text chunks (with [IMG_xxx] markers inlined)
        image_metas:      list of image metadata dicts
        chunk_image_ids:  per-chunk list of referenced image IDs
    """
    from docx import Document
    from docx.oxml.ns import qn

    try:
        doc = Document(file_path)
    except Exception as e:
        return [f"[DOCX解析错误] {str(e)}"], [], [[]]

    # ── Build rId → (bytes, ext) map from the docx zip relationships ──
    rid_to_image: Dict[str, Tuple[bytes, str]] = {}
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            try:
                rels_xml = zf.read("word/_rels/document.xml.rels").decode("utf-8")
                import xml.etree.ElementTree as ET
                ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
                root = ET.fromstring(rels_xml)
                for rel in root.findall("r:Relationship", ns):
                    target = rel.get("Target", "")
                    if "media/" not in target:
                        continue
                    rid = rel.get("Id")
                    if not rid:
                        continue
                    # Resolve relative path inside the zip
                    if target.startswith("/"):
                        full = target.lstrip("/")
                    elif target.startswith("word/"):
                        full = target
                    else:
                        full = "word/" + target
                    try:
                        ext = os.path.splitext(target)[1].lower()
                        if ext in IMAGE_EXTENSIONS:
                            rid_to_image[rid] = (zf.read(full), ext)
                    except KeyError:
                        pass
            except KeyError:
                pass
    except Exception as e:
        logger.warning("DOCX关系/图片提取失败: %s", e)

    # ── Walk body in document order ──
    image_metas: List[Dict[str, Any]] = []
    paragraphs: List[str] = []   # pre-chunking text segments (paragraphs joined linearly)
    table_chunks: List[str] = [] # each table becomes its own chunk
    seq_num = 0
    BLIP_EMBED_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"

    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag
        if tag == qn("w:p"):
            # Extract text via python-docx's Paragraph for accuracy
            from docx.text.paragraph import Paragraph
            para = Paragraph(child, doc)
            text = (para.text or "").strip()

            # Find all inline image references (a:blip elements)
            ids_here: List[str] = []
            for el in child.iter():
                if not el.tag.endswith("}blip"):
                    continue
                rid = el.get(BLIP_EMBED_ATTR)
                if not rid or rid not in rid_to_image:
                    continue
                img_bytes, img_ext = rid_to_image[rid]
                seq_num += 1
                desc = await _describe_image_with_vision(img_bytes, img_ext, vision_provider)
                img_id = _make_image_id(source_doc, seq_num, desc)
                local_path = _save_image(img_bytes, kb_id, doc_id, seq_num, img_ext)
                image_metas.append({
                    "id": img_id,
                    "doc_id": doc_id,
                    "kb_id": kb_id,
                    "seq_num": seq_num,
                    "page_num": 0,
                    "description": desc,
                    "local_path": local_path,
                    "source_doc": source_doc,
                })
                ids_here.append(img_id)
                text = (text + f" [{img_id}]").strip()

            if text:
                paragraphs.append(text)
        elif tag == qn("w:tbl"):
            from docx.table import Table
            tbl = Table(child, doc)
            md = _table_to_markdown(tbl)
            if md.strip():
                table_chunks.append(md)

    # ── Chunk linearized paragraph text (markers preserved by splitter) ──
    full_text = "\n".join(paragraphs)
    text_chunks = _split_text(full_text, chunk_size, chunk_overlap) if full_text.strip() else []

    # Per text chunk, harvest [IMG_xxx] markers
    chunk_image_ids: List[List[str]] = []
    for ch in text_chunks:
        ids = re.findall(r"\[IMG_[^\]]+\]", ch)
        chunk_image_ids.append([i.strip("[]") for i in ids])

    # ── Append each table as a standalone chunk (no images) ──
    for tc in table_chunks:
        text_chunks.append(tc)
        chunk_image_ids.append([])

    # Safety: if document had only images (no text/tables), still emit one chunk
    # per image so they remain searchable.
    if not text_chunks and image_metas:
        for meta in image_metas:
            text_chunks.append(f"[{meta['id']}] {meta['description']}")
            chunk_image_ids.append([meta["id"]])

    return text_chunks, image_metas, chunk_image_ids


async def _parse_pdf_with_media(
    file_path: str,
    doc_id: str,
    kb_id: str,
    chunk_size: int,
    chunk_overlap: int,
    vision_provider,
    source_doc: str,
) -> Tuple[List[str], List[Dict[str, Any]], List[List[str]]]:
    """Parse PDF using pymupdf with position-aware image insertion.

    For each page, text blocks and image bboxes are sorted by vertical position
    so the resulting linear text contains ``[IMG_xxx]`` markers at roughly the
    spot where the image appears in the page.
    """
    image_metas: List[Dict[str, Any]] = []

    # Try pymupdf for combined text + image extraction
    try:
        import fitz  # pymupdf
    except ImportError:
        # Fallback: plain text extraction, no images
        chunks = parse_pdf(file_path, chunk_size, chunk_overlap)
        return chunks, [], [[] for _ in chunks]

    try:
        pdf = fitz.open(file_path)
    except Exception as e:
        return [f"[PDF解析错误] {str(e)}"], [], [[]]

    page_texts: List[str] = []
    seq_num = 0

    try:
        for page_num, page in enumerate(pdf):
            # Collect items: (y_top, kind, payload)
            items: List[Tuple[float, str, Any]] = []

            # Text blocks
            try:
                blocks = page.get_text("blocks") or []
            except Exception:
                blocks = []
            for b in blocks:
                # block tuple: (x0, y0, x1, y1, text, block_no, block_type)
                if len(b) < 5:
                    continue
                btype = b[6] if len(b) >= 7 else 0
                if btype != 0:
                    continue  # only text blocks
                txt = (b[4] or "").strip()
                if txt:
                    items.append((float(b[1]), "text", txt))

            # Image blocks (with bbox)
            try:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    bboxes = page.get_image_rects(xref) or []
                    base_image = pdf.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = "." + base_image.get("ext", "png").lower()
                    if img_ext not in IMAGE_EXTENSIONS:
                        continue
                    y_top = float(bboxes[0].y0) if bboxes else 1e9  # unknown → end
                    items.append((y_top, "image", (img_bytes, img_ext)))
            except Exception as e:
                logger.warning("PDF page%d 图片提取失败: %s", page_num + 1, e)

            # Sort by vertical position
            items.sort(key=lambda x: x[0])

            page_buf: List[str] = []
            for _y, kind, payload in items:
                if kind == "text":
                    page_buf.append(payload)
                else:
                    img_bytes, img_ext = payload
                    seq_num += 1
                    desc = await _describe_image_with_vision(img_bytes, img_ext, vision_provider)
                    img_id = _make_image_id(source_doc, seq_num, desc)
                    local_path = _save_image(img_bytes, kb_id, doc_id, seq_num, img_ext)
                    image_metas.append({
                        "id": img_id,
                        "doc_id": doc_id,
                        "kb_id": kb_id,
                        "seq_num": seq_num,
                        "page_num": page_num + 1,
                        "description": desc,
                        "local_path": local_path,
                        "source_doc": source_doc,
                    })
                    page_buf.append(f"[{img_id}]")

            if page_buf:
                page_texts.append("\n".join(page_buf))
    finally:
        pdf.close()

    full_text = "\n\n".join(page_texts)
    text_chunks = _split_text(full_text, chunk_size, chunk_overlap) if full_text.strip() else []

    chunk_image_ids: List[List[str]] = []
    for ch in text_chunks:
        ids = re.findall(r"\[IMG_[^\]]+\]", ch)
        chunk_image_ids.append([i.strip("[]") for i in ids])

    # Safety: image-only PDF
    if not text_chunks and image_metas:
        for meta in image_metas:
            text_chunks.append(f"[{meta['id']}] {meta['description']}")
            chunk_image_ids.append([meta["id"]])

    return text_chunks, image_metas, chunk_image_ids


async def parse_document_with_media(
    file_path: str,
    doc_id: str,
    kb_id: str,
    chunk_size: int = 500,
    chunk_overlap: int = 50,
    vision_provider=None,
    source_doc: str = "",
) -> Tuple[List[str], List[Dict[str, Any]], List[List[str]]]:
    """Parse a document and extract images/tables.

    Returns ``(chunks, image_metas, chunk_image_ids)``:
        chunks            — text chunks (with [IMG_xxx] markers inlined when present)
        image_metas       — list of image metadata dicts
        chunk_image_ids   — per-chunk list of referenced image IDs (parallel to chunks)

    Falls back to plain parse_document for formats without media support, in
    which case chunk_image_ids is a list of empty lists.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".docx":
        return await _parse_docx_with_media(
            file_path, doc_id, kb_id, chunk_size, chunk_overlap, vision_provider, source_doc
        )
    elif ext == ".pdf":
        return await _parse_pdf_with_media(
            file_path, doc_id, kb_id, chunk_size, chunk_overlap, vision_provider, source_doc
        )
    else:
        # For other formats, just return plain chunks and no images
        chunks = parse_document(file_path, chunk_size, chunk_overlap)
        return chunks, [], [[] for _ in chunks]
